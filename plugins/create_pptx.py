"""
plugins/create_pptx.py — Generate PowerPoint (.pptx) presentations

Agent calls create_pptx with title + slides → gets a formatted .pptx file.
Supports: title slides, content slides with bullets, two-column layouts.

Usage by LLM:
    create_pptx({
        "filename": "presentation.pptx",
        "title": "Расчёт насыщения ТТ по ГОСТ",
        "author": "Pipeline Agent",
        "slides": [
            {"title": "Введение", "bullets": ["Цель работы", "Методология ГОСТ Р 58669-2019"]},
            {"title": "Формулы", "text": "Ψ = I₁ × (R₂ + R_load) × (1 + ω×T_a)"},
            {"title": "Результаты", "bullets": ["Время насыщения: 45мс", "Запас: 1.8x"], "notes": "Подробности в отчёте"},
            {"title": "Выводы", "bullets": ["ТТ удовлетворяет требованиям", "Рекомендации по замене"]}
        ]
    })
    → "Created presentation.pptx (5 slides incl. title, 34.2 KB)"

Requires: pip install python-pptx
"""

TOOL_DEFINITION = {
    "name": "create_pptx",
    "description": (
        "Create a PowerPoint (.pptx) presentation. "
        "Pass title + array of slides. Each slide can have: "
        "title, bullets (array of strings), text (single block), "
        "and notes (speaker notes). First slide is auto-generated as title slide."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Output filename, e.g. 'presentation.pptx'",
            },
            "title": {
                "type": "string",
                "description": "Presentation title (shown on the title slide)",
            },
            "author": {
                "type": "string",
                "description": "Author/subtitle on title slide (optional)",
            },
            "slides": {
                "type": "string",
                "description": (
                    "JSON array of slides. Each: "
                    "{\"title\": \"...\", \"bullets\": [\"...\", \"...\"], \"text\": \"...\", \"notes\": \"...\"}. "
                    "Use 'bullets' for bullet-point slides, 'text' for single-paragraph slides. "
                    "Example: [{\"title\":\"Goals\",\"bullets\":[\"Item 1\",\"Item 2\"]},{\"title\":\"Details\",\"text\":\"Full paragraph...\"}]"
                ),
            },
        },
        "required": ["filename", "title", "slides"],
    },
}

AGENTS = ["coder", "writer"]


def execute(args: dict) -> str:
    import json
    from pathlib import Path

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
    except ImportError:
        return "ERROR: python-pptx not installed. Run: pip install python-pptx"

    filename = args.get("filename", "presentation.pptx")
    title    = args.get("title", "Untitled Presentation")
    author   = args.get("author", "")
    slides_raw = args.get("slides", "[]")

    # Parse slides
    try:
        if isinstance(slides_raw, str):
            slides = json.loads(slides_raw)
        else:
            slides = slides_raw
    except json.JSONDecodeError as e:
        return f"ERROR: Invalid JSON in slides: {e}"

    if not isinstance(slides, list):
        return "ERROR: slides must be a JSON array"

    prs = Presentation()

    # Set slide dimensions (16:9)
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title slide ───────────────────────────────────────────────────────────
    layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout_title)

    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = author or ""

    # ── Content slides ────────────────────────────────────────────────────────
    layout_content = prs.slide_layouts[1]  # Title and Content layout
    layout_blank   = prs.slide_layouts[6]  # Blank layout

    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            continue

        slide_title  = s.get("title", f"Slide {i+1}")
        bullets      = s.get("bullets", [])
        text         = s.get("text", "")
        notes        = s.get("notes", "")

        if bullets:
            # Bullet point slide
            slide = prs.slides.add_slide(layout_content)
            slide.shapes.title.text = slide_title

            body = slide.placeholders[1]
            tf   = body.text_frame
            tf.clear()

            for j, bullet in enumerate(bullets if isinstance(bullets, list) else [bullets]):
                if j == 0:
                    tf.text = str(bullet)
                else:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0

            # Style bullets
            for para in tf.paragraphs:
                para.font.size = Pt(18)

        elif text:
            # Text-only slide
            slide = prs.slides.add_slide(layout_content)
            slide.shapes.title.text = slide_title

            body = slide.placeholders[1]
            body.text_frame.clear()
            body.text_frame.text = str(text)

            for para in body.text_frame.paragraphs:
                para.font.size = Pt(16)

        else:
            # Title-only slide
            slide = prs.slides.add_slide(layout_content)
            slide.shapes.title.text = slide_title

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = str(notes)

    # Save
    output = Path("workspace") / filename
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    size = output.stat().st_size
    size_str = f"{size/1024:.1f} KB" if size > 1024 else f"{size} bytes"
    total_slides = len(slides) + 1  # +1 for title slide
    return f"Created {filename} ({total_slides} slides incl. title, {size_str})"
